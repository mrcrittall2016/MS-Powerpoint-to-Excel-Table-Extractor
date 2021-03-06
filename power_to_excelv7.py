# This code is designed to extract text from a TABLE shape type in a powerpoint slide and import it to Microsoft Excel. 
# It's first task is to identify how many different TYPES of table there are - meaning how many different heading lists there are
# Try to simplify code so only goes through presentation once... also add command-line arguments
# Now attempt to turn program into class - good stack overflow article on this: https://stackoverflow.com/questions/28684053/python-turn-giant-function-into-class
# Create Child class of program that performs same function as parent but instead also combines all PK data into single spreadsheet
# Copy data across to master combined table by row

import sys #import argv
import os
from pptx import Presentation
import openpyxl
from openpyxl.styles import Alignment, Font


class Powerpoint_to_Excel(object):
	
	# Instantiate object. Essentially the below can be thought of akin to global variables in normal functional program
	def __init__(self, powerpoint_file):
		self.prs = Presentation(powerpoint_file)
		self.wb = openpyxl.load_workbook('example.xlsx')
		self.header_list = []
		self.header_templates = {}
		self.template_tracker = {}
		self.template_index = 1
		self.total_rows = 0
		self.number_of_tables = 0
		self.new_templates = 0
		sheet = self.wb.active
		sheet.title = "Template " + str(self.template_index)
	
	
	#Fill any blank headers with "Blank 1,2 etc.."
	def fill_in_blanks(self, header_list):
		# First cycle through header_list and replace blanks with "Blank 1.." etc.
		blank_index = 1		
	
		#Replace any blank headers
		for i, entry in enumerate(header_list):
			#print "The current list entry is: %s" % entry		
			if entry == '':	
				#print "Nothing here..."			
				header_list[i] = "Blank " + str(blank_index)
				blank_index += 1
	
	#Function for calculating total number of rows in presentation
	def row_calculator(self, prs):
		
		slides = self.prs.slides
		
		# Count number of rows
		rows = 0
	
		for slide in slides:
			for shape in slide.shapes:
				if shape.shape_type == 19:
					table = shape.table
					for row_count, row in enumerate(table.rows, 0):
						if row_count != 0:
							rows += 1
		return rows
	
	
	#Function to manipulate Excel workbook once have identified correct table-type or template
	def add_data(self, index, list, table, workbook):
		#Create sheet for template if not already exists
		if index not in workbook.sheetnames:
			#Create new sheet in Excel with template name if does not already exist
			workbook.create_sheet(index)						

		# Now that we have an Excel workbook with a sheet per template, we need to populate it. Try putting in just headers first
		# Put column_headers into Excel and style. Only really want to do this once though... and probably have as separate function
		for column_count, header in enumerate(list, 1):					
			#print "\nThe list index here is: %s\n" % list_template_index
	
			#print "\nThe header is: %s\n" % header
	
			sheet_selected = workbook[index]	
			cell_pos = sheet_selected.cell(row=1, column=column_count)
			cell_pos.value = header	
			# Ensures data is centered within cell
			cell_pos.alignment = Alignment(horizontal="center", wrapText=True)						
			cell_pos.font = Font(bold=True)



		#Now that headers are filled out in each template...need to add in data
		for row_index, row in enumerate(table.rows):						
			single_row_empty = [table.cell(row_index, column).text_frame.text for column in range(0, len(list)) if table.cell(row_index, column).text_frame.text == '' or table.cell(row_index, column).text_frame.text == None]
			#print "\nRow %s in table is: %s\n" % (row_index, single_row)

			#print "\nThe length of a single row is: %s\n" % len(single_row)
			
			# If row is empty.. jump over this row
			if len(single_row_empty) == len(list):
				continue
			
			else:
				
				#Re-set column#
				column = 0
				
				single_row = [table.cell(row_index, column).text_frame.text for column in range(0, len(list))]
				
				#Zip up row with header template only if same length
				if row_index != 0 and len(single_row) == len(list):							
	
					#print "\nLength of single row is : %s\n" % len(single_row) #17
	
					#print "\nLength of header_list row is : %s\n" % len(header_list) #17	 													
	
					zipped_row = dict(zip(list, single_row))
	
	
					# Check zipped row
					#for key, value in zipped_row.items():
						#print "\n%s : %s\n" % (key, value)
	
					#print "Zipped row is: %s" % zipped_row #16
	
					#Zipped row changes length from 17 to 16... why? Because two headers of same name i.e. blanks... so need to name blanks beforehand
					#print "\nLength of zipped row is : %s\n" % len(zipped_row)							
	
					#Now send to relevant Excel sheet/template
					sheet_to_insert_into = workbook[index]  #Also try sheet_selected
	
					#How do we know how many rows to insert for this table?
					#rows_to_insert = row_calculator(table)
	
					self.Excel_add(sheet_to_insert_into, zipped_row)



	# Function for adding row/dictionary into Excel
	def Excel_add(self, sheet, row_to_insert):
	
		entries = 0	
	
		#print "Max number of rows is: ", row_number	
	
		#print "The length of the row to insert is: %s" % len(row_to_insert)
	
		#Go through spreadsheet. Need to ensure provide enough rows
		for rowNum in range (1, self.total_rows):
			for colNum in range (1, len(row_to_insert) + 1):
			
				# Entries variable ensures do not get stuck in this excel loop with same dictionary ie only
				# adds in each row once.
					
				if entries == len(row_to_insert):
					return
			
				cell_val = sheet.cell(row=rowNum, column=colNum).value
			
				#print "Cell val is: ", cell_val
			
				if cell_val == None:
				
					#print "Cell is empty!"
				
					# Check header of that column against a key in the dictionary
					col_value = sheet.cell(row=1, column=colNum).value
				
					#print "\nheader_value is: \n", col_value
					#sheet.cell(row=rowNum, column=colNum).value = table_count
				
					for key, value in row_to_insert.items():
						if col_value == key:
							entries += 1
							sheet.cell(row=rowNum, column=colNum).value = row_to_insert[key]
						
							# Ensures data is centered within cell
							sheet.cell(row=rowNum, column=colNum).alignment = Alignment(horizontal="center", wrapText=True)
				

	# Function to add in new table heading templates to dictionary
	def add_to_dictionary(self, list, dictionary, slide):	
	
		unique_entry_counter = 0
	
		for list_template_index, list_already_there in dictionary.items():
		
			#print "List already there is: %s and its length is %s" % (list_already_there, len(list_already_there))	
			#print "\nList to add is %s and its legnth is: %s\n" % (list, len(list))			
		
			# If template is unique
			if list != list_already_there:						
				unique_entry_counter = unique_entry_counter + 1
			
			# If template is not unique, which template is it equal to? Return the template index so no which sheet to look at in Excel workbook
			elif list == list_already_there:
				return list_template_index
						
	
		# If the list to add is different to ALL of the dictionary entries (ie the length of the dictionary), then add it in		
		if unique_entry_counter == len(dictionary):		
			#print "print uniqueness is: %d" % unique_entry_counter
			#print "\nAdding new template\n"		
			
			#print "List now is: %s" % list		
		
			#Which slide was the unique table found on and what table number? Also what was its length?					
			template_length = len(list)			
			
			self.template_tracker["Template " + str(self.template_index)] = str(self.number_of_tables) + " on slide " + str(slide)	
			self.template_tracker["Template " + str(self.template_index) + " length: "] = template_length	
		
			#global new_templates 
			self.new_templates += 1
		
			#Keep track of old template name before updating index
			template_added = "Template " + str(self.template_index)
		
			# Add new template to dictionary
			dictionary[template_added] = list				
		
			#global template_index
			self.template_index += 1
		
			return template_added
	


	# Go through presentation analysing different table types and creating a new Excel sheet of stored data for each type
	def extract_and_transfer(self, workbook):	
		for slide_number, slide in enumerate(self.prs.slides, 1):
			for shape in slide.shapes:
				if shape.shape_type == 19:
					#global number_of_tables
					self.number_of_tables += 1
			
					#A table. Have table.columns, table.rows and table.cell(row, column).text_frame.text for cell value
					table = shape.table
					number_of_columns = 0
			
					# Now we want to get just the headers which should be in row 0. But how do we know how many columns we have? See below
					for column_number, column in enumerate(table.columns, 1):				
						number_of_columns = column_number
			
					#print "On slide %s there are %s columns" % (slide_number, number_of_columns)
			
					#Based on the number of columns we have just calculated, create a list of headers
					#header_list = [table.cell(0, column).text_frame.text for column in range(0, number_of_columns)]
					header_list = []					
					
					# Really need to re-name headers appropriately... ie Gender or Fed/Fasted
					for column in range(0, number_of_columns):						
						
						# If no row beneath apparent header, carry on
						try:
							if table.cell(1, column).text_frame.text == 'F' or table.cell(1, column).text_frame.text == 'M':
								#header_list.append('Gender')
								table.cell(0, column).text_frame.text = 'Gender'
							if table.cell(1, column).text_frame.text == 'Fed' or table.cell(1, column).text_frame.text == 'Fasted':
								#header_list.append('Fed/Fasted')
								table.cell(0, column).text_frame.text = 'Fed/Fasted'
						except:
							continue
						
						header_list.append(table.cell(0, column).text_frame.text)
					
					
					#print "\nChecking header_list: %s\n" % header_list
				
					# Now want to add list to the header_template dictionary only if dictionary is empty or if different to a header template already located in the dictionary
			
					# Is dictionary already empty? If so add header_list
					if not self.header_templates:					
					
						list_template_index = self.add_to_dictionary(header_list, self.header_templates, slide_number)	
					
						# Now based on returned template index, find correct sheet in Excel and add data	
						#print "\nReturn is: %s\n" % list_template_index	
					
						# Pass list_template_index to separate function in future, with header_list and table					
						self.add_data(list_template_index, header_list, table, self.wb)
			
					# If it is not empty...			
					elif self.header_templates:
			
						#print "\nDictionary is not empty on slide %s\n" % slide_number				
			
						#Pass to function to add to dictionary if not already there
						list_template_index = self.add_to_dictionary(header_list, self.header_templates, slide_number)					
				
						# Now based on returned template index, find correct sheet in Excel and add data
						#print "\nReturn is: %s\n" % list_template_index
				
					
						# Pass list_template_index to separate function in future, with header_list and table					
						self.add_data(list_template_index, header_list, table, self.wb)
	
	
	
	

	def main(self):
		
		#Call member function as described here: https://stackoverflow.com/questions/5615648/python-call-function-within-class
		#print "The number of rows is: %s" % self.row_calculator(self.prs)		
		
		self.total_rows = self.row_calculator(self.prs)
		
		self.extract_and_transfer(self.wb)
		
		self.wb.save('example_copy_OOP.xlsx')
		
		print "\nProcess completed successfully\n"



# Inherit all methods and functionality from Powerpoint_to_Excel class above. Specialise this class to have combine method
class Combine(Powerpoint_to_Excel):		
	
	def combine(self):	
		
		length_list = {}
		
		for key, value in self.template_tracker.items():
			#print "\n%s : %s\n" % (key, value)
			
			#Extract into its own dictionary
			if "length" in key:
				length_list[key] = value
			
		
		# From new dictionary, identify Template with highest value
		#print length_list
		
		# Obtain maximum value in dict
		#print max(length_list.values())
		
		#Obtain key associated with max value. See here: https://stackoverflow.com/questions/26871866/print-highest-value-in-dict-with-key
		maximum = max(length_list, key=length_list.get)
		#print maximum
		
		# Strip out integer from key...
		maximum = "Template " + ''.join([s for s in maximum if s.isdigit()])
				
		# Now with longest template identified, extract headers from header_templates dictionary
		#print self.header_templates[maximum]
		
		# Create this sheet in Excel Workbook
		self.wb.create_sheet("Combined")
		
		for column_count, header in enumerate(self.header_templates[maximum], 1):	
			sheet_selected = self.wb["Combined"]	
			cell_pos = sheet_selected.cell(row=1, column=column_count)
			cell_pos.value = header	
			# Ensures data is centered within cell
			cell_pos.alignment = Alignment(horizontal="center", wrapText=True)						
			cell_pos.font = Font(bold=True)
			
		# Go through each sheet and take out a row at a time
		for index in range(1, len(self.header_templates) + 1):
		
			#print index
			index = "Template " + str(index)
		
			# Activate Excel sheet
			active_sheet = self.wb[index]
			
			#Go through activated sheet - assume first row is header row. 
			for rowNum in range (2, self.total_rows):
				
				single_row_empty = [active_sheet.cell(row=rowNum, column=colNum).value for colNum in range (1, len(self.header_templates[index]) + 1) if active_sheet.cell(row=rowNum, column=colNum).value == None or active_sheet.cell(row=rowNum, column=colNum).value == '']
				
				#print "\nsingle_row_empty is: %s\n" % single_row_empty
				#print "\nCurrent rowNum is: %s\n" % rowNum
				#print "\nCurrent sheet is: %s\n" % index
				
				# If single_row is completely empty, break
				if len(single_row_empty) == len(self.header_templates[index]):
					break
				
				# Go back through row and store if not empty row
				else:
					
					# Re-set colNum
					colNum = 1
					
					single_row = [active_sheet.cell(row=rowNum, column=colNum).value for colNum in range (1, len(self.header_templates[index]) + 1)]
					 
		
					# Now zip up single row with header_template for that sheet
					zipped_row = dict(zip(self.header_templates[index], single_row))
				
					#print "\nzipped_row is: %s\n" % zipped_row
				
					# Now go through combined sheet
					combined_sheet = self.wb["Combined"]
				
					for row_num in range(2, self.total_rows):
					
						# Check to see if row is completely empty..
						single_row_empty = [combined_sheet.cell(row=row_num, column=col_num).value for col_num in range (1, len(self.header_templates[maximum]) + 1) if combined_sheet.cell(row=row_num, column=col_num).value == None]
									
						
						
						#print "Length os single_row is: %s" % len(single_row)						
					
						if len(single_row_empty) == len(self.header_templates[maximum]):
							#print "Empty row found!"
							#print "Value of col_num is: %s" % col_num
							#print "Value of row_num is: %s" % row_num
						
							# Need to re-set column variable back to 1 and go through row and compare against zipped_row
							col_num = 1
							#print "\nCurrent rowNum is: %s\n" % row_num
							#print "\nzipped_row to insert is: %s\n" % zipped_row
							#print "\nCurrent sheet is: %s\n" % index
							
							# Now using re-set column number, iterate through from beginning of row and place in correct value from zipped_row
							for col_num in range(1, len(self.header_templates[maximum]) + 1):
								#print "\nInside new loop\n"
								#print "Value of col_num is now: %s" % col_num
								#print "Value of row_num is now: %s" % row_num
								for key, value in zipped_row.items():
									#print "Value of zipped row dict is: %s" % zipped_row[key]
									#print "Value of combined sheet header is: %s" % combined_sheet.cell(row=1, column=col_num).value
									if key == combined_sheet.cell(row=1, column=col_num).value:
										combined_sheet.cell(row=row_num, column=col_num).value = zipped_row[key]
						
							break			
		
	
	# Override main method so that save as different file
	def main(self):
		self.total_rows = self.row_calculator(self.prs)		
		
		self.extract_and_transfer(self.wb)		
		
		self.combine()
		
		self.wb.save('example_copy_OOPv4.xlsx')	
		
		print "\nThere should be %s rows in this Presentation\n" % self.total_rows
		
		print "\nProcess completed successfully\n"
		
		# Now need to figure out what is longest list... ie list with most column headings. 
		# Probably want to modify add_to_dictionary method to measure length of each "unique" template it stores
		
		
		


#Get command line arguments
#script, input = argv

# Ensure provide just 1 command-line argument
if len(sys.argv) > 2:
	print "\nPlease provide only one command-line argument. Thankyou\n"
	sys.exit(1)
	
try:
	input = sys.argv[1]
	
except IndexError:
	print "Usage: myprogram.py <arg1>"
	print "Please provide a MS Powerpoint file name for analysis"
	sys.exit(1)

# Check that file provided is reachable
if os.path.isfile(input):
	#Instantiate object and run 'main' method
	test_run = Combine(input)
	test_run.main() 

else:
	print "\nSorry, file does not exist. Try providing an alternate file-path\n"

# load a presentation
#path_to_pptx_file = "PK_Feb_test.pptx"
#path_to_pptx_file = "PK_Feb.pptx"
#path_to_pptx_file = "PK_Feb_test_1.pptx"














