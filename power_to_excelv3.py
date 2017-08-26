# This code is designed to extract text from a TABLE shape type in a powerpoint slide and import it to Microsoft Excel. 
# It's first task is to identify how many different TYPES of table there are - meaning how many different heading lists there are
from pptx import Presentation
import openpyxl
from openpyxl.styles import Alignment, Font

# load a presentation
path_to_pptx_file = "PK_Feb_test.pptx"
#path_to_pptx_file = "PK_Feb.pptx"
#path_to_pptx_file = "PK_Feb_test_1.pptx"

prs = Presentation(path_to_pptx_file)

number_of_tables = 0

header_list = []

header_templates = {}
template_tracker = {}
template_index = 1

new_templates = 0
templates_matched = 0

total_rows = 0

#Fill any blank headers with "Blank 1,2 etc.."
def fill_in_blanks(header_list):
	# First cycle through header_list and replace blanks with "Blank 1.." etc.
	blank_index = 1		
	
	#Replace any blank headers
	for i, entry in enumerate(header_list):
		#print "The current list entry is: %s" % entry		
		if entry == '':	
			#print "Nothing here..."			
			header_list[i] = "Blank " + str(blank_index)
			blank_index += 1
	
#Function for calculating total number of rows in presentation
def row_calculator(prs):
	slides = prs.slides
	# Count number of rows
	rows = 0
	
	for slide in slides:
		for shape in slide.shapes:
			if shape.shape_type == 19:
				table = shape.table
				for row_count, row in enumerate(table.rows, 0):
					if row_count != 0:
						rows += 1
	return rows


# Function for adding row/dictionary into Excel
def Excel_add(sheet, row_to_insert):
	
	entries = 0	
	
	#print "Max number of rows is: ", row_number	
	
	#print "The length of the row to insert is: %s" % len(row_to_insert)
	
	#Go through spreadsheet. Need to ensure provide enough rows
	for rowNum in range (1, total_rows):
		for colNum in range (1, len(row_to_insert) + 1):
			
			# Entries variable ensures do not get stuck in this excel loop with same dictionary ie only
			# adds in each row once.
					
			if entries == len(row_to_insert):
				return
			
			cell_val = sheet.cell(row=rowNum, column=colNum).value
			
			#print "Cell val is: ", cell_val
			
			if cell_val == None:
				
				#print "Cell is empty!"
				
				# Check header of that column against a key in the dictionary
				col_value = sheet.cell(row=1, column=colNum).value
				
				#print "\nheader_value is: \n", col_value
				#sheet.cell(row=rowNum, column=colNum).value = table_count
				
				for key, value in row_to_insert.items():
					if col_value == key:
						entries += 1
						sheet.cell(row=rowNum, column=colNum).value = row_to_insert[key]
						
						# Ensures data is centered within cell
						sheet.cell(row=rowNum, column=colNum).alignment = Alignment(horizontal="center", wrapText=True)
				

# Function to add in new table heading templates to dictionary
def add_to_dictionary(list, dictionary, slide):	
	
	unique_entry_counter = 0
	
	for list_template_index, list_already_there in dictionary.items():
		
		#print "List already there is: %s and its length is %s" % (list_already_there, len(list_already_there))	
		#print "\nList to add is %s and its legnth is: %s\n" % (list, len(list))			
		
		# If template is unique
		if list != list_already_there:						
			unique_entry_counter = unique_entry_counter + 1				
	
	# If the list to add is different to ALL of the dictionary entries (ie the length of the dictionary), then add it in		
	if unique_entry_counter == len(dictionary):		
		#print "print uniqueness is: %d" % unique_entry_counter
		#print "\nAdding new template\n"		
			
		#print "List now is: %s" % list		
		
		#Which slide was the unique table found on and what table number?
		template_tracker["Template " + str(template_index)] = str(number_of_tables) + " on slide " + str(slide)		
		
		global new_templates 
		new_templates = new_templates + 1
		
		dictionary["Template " + str(template_index)] = list
		global template_index
		template_index = template_index + 1	
		


# First go through presentation and analyse all different table types
def extract_table_headers():
	
	for slide_number, slide in enumerate(prs.slides, 1):
		for shape in slide.shapes:
			if shape.shape_type == 19:
				global number_of_tables
				number_of_tables = number_of_tables + 1
			
				#A table. Have table.columns, table.rows and table.cell(row, column).text_frame.text for cell value
				table = shape.table
				number_of_columns = 0
			
				# Now we want to get just the headers which should be in row 0. But how do we know how many columns we have? See below
				for column_number, column in enumerate(table.columns, 1):				
					number_of_columns = column_number
			
				#print "On slide %s there are %s columns" % (slide_number, number_of_columns)
			
				#Based on the number of columns we have just calculated, create a list of headers
				header_list = [table.cell(0, column).text_frame.text for column in range(0, number_of_columns)]
			
				#print header_list				
				
				fill_in_blanks(header_list)				
			
				# Now want to add list to the header_template dictionary only if dictionary is empty or if different to a header template already located in the dictionary
			
				# Is dictionary already empty? If so add header_list
				if not header_templates:
					
					'''
					global template_index
					header_templates["Template " + str(template_index)] = header_list
					template_index = template_index + 1		
					'''
					add_to_dictionary(header_list, header_templates, slide_number)			
					
			
				# If it is not empty...			
				elif header_templates:
				
					#print "\nDictionary is not empty on slide %s\n" % slide_number				
				
					#Pass to function to add to dictionary if not already there
					add_to_dictionary(header_list, header_templates, slide_number)					




extract_table_headers()

#How many total rows are there in the presentation
total_rows = row_calculator(prs)

print "\nThere are %s total rows in this presentation\n" % total_rows

#Re-set template_index
template_index = 1

#Ok, now that we have our dictionary of templates, we want to go through presentation again and match each table to a template.
#If we find a match, then we can create a new sheet in the Excel file dedicated to that table type

# Let's ensure the Excel file is open first
wb = openpyxl.load_workbook('example.xlsx')
sheet = wb.active
sheet.title = "Template " + str(template_index)



for slide_number, slide in enumerate(prs.slides):
	for shape in slide.shapes:
		if shape.shape_type == 19:
			
			table = shape.table
			
			#How many rows in this table?
			#rows = row_calculator(table)
						
			number_of_columns = 0
			
			# Now we want to get just the headers which should be in row 0. But how do we know how many columns we have? See below
			for column_number, column in enumerate(table.columns, 1):				
				number_of_columns = column_number
			
			#print "number_of_columns is: %s" % number_of_columns
			
			header_list = [table.cell(0, column).text_frame.text for column in range(0, number_of_columns)]			
			
			fill_in_blanks(header_list)	
			
			# Now compare header_list against what is in dictionary of templates
			for list_template_index, list_to_compare_against in header_templates.items():
				
				
				#print "\nThe length of the dictionary is: %d\n" % len(header_templates)				
				#print "\nThe length of header_list is: %s\n" % len(header_list)				
				#print "\nThe length of dest_list is: %s\n" % len(list_to_compare_against)
				
				
				#If match header to one of unique templates in dictionary...
				if header_list == list_to_compare_against:
					#print "\nFound match!\n"
					templates_matched = templates_matched + 1					
					# Which template did we find?
					#print "\nWe found template: %s\n" % list_template_index
					
					#Create sheet for template if not already exists
					if list_template_index not in wb.sheetnames:					
						
						#Create new sheet in Excel with template name if does not already exist
						wb.create_sheet(list_template_index)						
						
					# Now that we have an Excel workbook with a sheet per template, we need to populate it. Try putting in just headers first
					# Put column_headers into Excel and style. Only really want to do this once though... and probably have as separate function
					for column_count, header in enumerate(header_list, 1):					
						#print "\nThe list index here is: %s\n" % list_template_index
						sheet_selected = wb[list_template_index]	
						cell_pos = sheet_selected.cell(row=1, column=column_count)
						cell_pos.value = header	
						# Ensures data is centered within cell
						cell_pos.alignment = Alignment(horizontal="center", wrapText=True)						
						cell_pos.font = Font(bold=True)
					
					
					
					#Now that headers are filled out in each template...need to add in data
					for row_index, row in enumerate(table.rows):						
						single_row = [table.cell(row_index, column).text_frame.text for column in range(0, number_of_columns)]
						#print "\nRow %s in table is: %s\n" % (row_index, single_row)
						
						#print "\nThe length of a single row is: %s\n" % len(single_row)
						
												
						#Zip up row with header template only if same length
						if row_index != 0 and len(single_row) == len(header_list):							
							
							#print "\nLength of single row is : %s\n" % len(single_row) #17
							
							#print "\nLength of header_list row is : %s\n" % len(header_list) #17	 													
							
							zipped_row = dict(zip(header_list, single_row))
							
							'''
							# Check zipped row
							for key, value in zipped_row.items():
								print "\n%s : %s\n" % (key, value)
							'''
							#print "Zipped row is: %s" % zipped_row #16
							
							#Zipped row changes length from 17 to 16... why? Because two headers of same name i.e. blanks... so need to name blanks beforehand
							#print "\nLength of zipped row is : %s\n" % len(zipped_row)							
							
							#Now send to relevant Excel sheet/template
							sheet_to_insert_into = wb[list_template_index]  #Also try sheet_selected
							
							#How do we know how many rows to insert for this table?
							#rows_to_insert = row_calculator(table)
							
							Excel_add(sheet_to_insert_into, zipped_row)
						

# Save Excel document 
wb.save('example_copy4.xlsx')


#print "\nTemplates matched is: %s\n" % templates_matched
		
#Should be ~248 tables... which is correct
#print "\nIn this presentation there are %s tables\n" % number_of_tables


#print "\nThe number of new templates is: %s" % new_templates

'''
#Check headers dictionary
for key, value in header_templates.items():
	print "%s is: %s" % (key, value)
'''	

'''
#Check template_tracker dictionary
for key, value in template_tracker.items():
	print "%s is: table %s" % (key, value)
'''
'''
#Dictionary sorted by key
for key in sorted(template_tracker.iterkeys()):
	print "%s is: table %s" % (key, template_tracker[key])
'''





















