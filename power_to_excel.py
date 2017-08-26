# This code is designed to extract text from a TABLE shape type in a powerpoint slide and import it to Microsoft Excel

from pptx import Presentation
import openpyxl
from openpyxl.styles import Alignment, Font

# load a presentation
path_to_pptx_file = "power_test2.pptx"

prs = Presentation(path_to_pptx_file)


# Function for checking how many rows of data need to extract from Powerpoint presentation
def row_calculator(prs):
	slides = prs.slides
	# Count number of rows
	rows = 0
	
	for slide in slides:
		for shape in slide.shapes:
			if shape.shape_type == 19:
				table = shape.table
				for row_count, row in enumerate(table.rows, 0):
					if row_count != 0:
						rows += 1
	return rows
	
	
	
# Function for adding row/dictionary into Excel
def Excel_add(sheet, DATA, column_headers, row_number):
	
	entries = 0
	
	print "Max number of rows is: ", row_number	
	
	for rowNum in range (1, row_number + 2):
		for colNum in range (1, len(column_headers) + 1):
			
			# Entries variable ensures do not get stuck in this excel loop with same dictionary ie only
			# adds in each row once.
			if entries == len(column_headers):
				return
			
			cell_val = sheet.cell(row=rowNum, column=colNum).value
			
			#print "Cell val is: ", cell_val
			
			if cell_val == None:
				
				#print "Cell is empty!"
				
				# Check header of that column against a key in the dictionary
				col_value = sheet.cell(row=1, column=colNum).value
				
				#print "header_value is: ", col_value
				
				for key in DATA:
					if col_value == key:
						entries += 1
						sheet.cell(row=rowNum, column=colNum).value = DATA[key]
						
						# Ensures data is centered within cell
						sheet.cell(row=rowNum, column=colNum).alignment = Alignment(horizontal="center", wrapText=True)
				

#Apparently table is a type of placeholder shape and has a unique index or id. See here:
#http://python-pptx.readthedocs.io/en/latest/user/placeholders-using.html

# The code below identifies what shape_types are embedded within the slide... for test should
# identify TABLE shape type. Just look at first slide to obtain table headings. 


# First slide
slide = prs.slides[0]

# Extract headers from first slide only
for shape in slide.shapes:
	
	# The idx value of a placeholder is the integer key of the slide layout placeholder 
	# it inherits properties from. TABLE has a key of 19 
	if shape.shape_type == 19:
		print "table found..."

		# Here shape represents a graphic_frame object of which table is an object within this
		table = shape.table

		# List for column headers
		column_headers = []

		# Enumerate - adds counter to for loop. See here for explanation: http://stackoverflow.com/questions/22171558/what-does-enumerate-mean  
		for row_count, row in enumerate(table.rows, 0):
			for column_count, column in enumerate(table.columns, 0):
				
				# Assume first row is table_headers
				if (row_count == 0):
					
					header = table.cell(row_count, column_count).text_frame.text

					# Store column_headers in list
					column_headers.append(header)

# Check how many rows of data to extract
row_number = row_calculator(prs)

# Now export headers to Excel... first load a workbook from current directory. See this website for a 
# tutorial on this: https://automatetheboringstuff.com/chapter12/ 
wb = openpyxl.load_workbook('example.xlsx')

sheet = wb.active

sheet.title = "All Data"

# Put column_headers into Excel and style
for column_count, header in enumerate(column_headers, 1):
	
	cell_pos = sheet.cell(row=1, column=column_count)
	cell_pos.value = header
	
	# Ensures data is centered within cell
	cell_pos.alignment = Alignment(horizontal="center", wrapText=True)
	#cell_pos.alignment = Alignment(vertical="justify")
	
	# Embolden font in header cells
	cell_pos.font = Font(bold=True)


# Now create a dictionary using the column headers as keys
DATA = {}

for index in range(0, len(column_headers)):
	print "Column header: ", column_headers[index]
	# Create dictionary
	DATA[column_headers[index]] = ''


print "Dictionary is: ", DATA


# Now that we have the dictionary structure...populate with correct values. Go through entire presentation this time
for slide in prs.slides:
	for shape in slide.shapes:
		if shape.shape_type == 19:
		
			print "Found table"
			key_value_pairs = 0
			table = shape.table
			
			# Found table, now iterate through columns and rows
			for row_count, row in enumerate(table.rows, 0):
				for column_count, column in enumerate(table.columns, 0):
					cell_value = table.cell(row_count, column_count).text_frame.text
					# Make sure not on first row ie header row
					if row_count != 0:
						# Check header value
						header_value = table.cell(0, column_count).text_frame.text
							
							
						# Iterate through dictionary keys. If key matches header value...
						for key in DATA:
						
							if header_value == key:
									
										
								DATA[key] = cell_value
								print DATA[key]
								key_value_pairs += 1
								print "key_value_pairs is: ", key_value_pairs
								
								
								if key_value_pairs == len(column_headers):
									
									# Reset pair variable
									key_value_pairs = 0
									
									# Call function to put row into Excel
									print "Key-value pairs complete: ", DATA
									
									# Now put into Excel if dictionary is complete
									Excel_add(sheet, DATA, column_headers, row_number)
									
									

# Save Excel document 
wb.save('example_copy.xlsx')







